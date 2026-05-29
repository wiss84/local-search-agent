"""
PPTX parser for the Local Search Agent ingestion pipeline.

Uses python-pptx for PowerPoint presentation extraction.

Strategy:
- Each slide becomes a Markdown section (## Slide N: <title>)
- Title, body text, and table content are extracted per slide
- Speaker notes are appended under each slide as a blockquote
- Empty slides are skipped
- Text order follows the slide's shape z-order (top-to-bottom)

Install: pip install "python-pptx>=0.6.23"
"""

from __future__ import annotations

import logging
import os
from typing import Optional

from local_search_agent.core.document_node import DocumentNode
from local_search_agent.ingestion.cleaner import clean
from local_search_agent.ingestion.parser import BaseParser, ParserError

logger = logging.getLogger(__name__)


def _table_to_markdown(table) -> str:
    """Convert a python-pptx Table object to a Markdown table string."""
    rows = []
    for row in table.rows:
        cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
        rows.append("| " + " | ".join(cells) + " |")

    if not rows:
        return ""

    # Insert separator after header row
    col_count = len(table.rows[0].cells)
    separator = "| " + " | ".join(["---"] * col_count) + " |"
    rows.insert(1, separator)
    return "\n".join(rows)


def _slide_to_markdown(slide, slide_num: int) -> str:
    """
    Convert a single python-pptx Slide to a Markdown section.

    Returns empty string if the slide has no extractable text content.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: PLC0415

    title_text = ""
    body_parts: list[str] = []
    table_parts: list[str] = []
    notes_text = ""

    # Extract title separately for the section heading
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        title_text = slide.shapes.title.text.strip()

    # Extract all shapes
    for shape in slide.shapes:
        # Skip the title shape — already handled above
        if shape == slide.shapes.title:
            continue

        # Tables
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            md_table = _table_to_markdown(shape.table)
            if md_table:
                table_parts.append(md_table)
            continue

        # Text frames
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                line = para.text.strip()
                if line:
                    body_parts.append(line)

    # Speaker notes
    try:
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_raw = notes_frame.text.strip()
            if notes_raw:
                notes_text = notes_raw
    except Exception:
        pass  # Notes extraction is best-effort

    # Build section
    has_content = title_text or body_parts or table_parts or notes_text
    if not has_content:
        return ""

    heading = f"## Slide {slide_num}: {title_text}" if title_text else f"## Slide {slide_num}"
    parts: list[str] = [heading]

    if body_parts:
        parts.append("\n".join(body_parts))

    if table_parts:
        parts.extend(table_parts)

    if notes_text:
        quoted = "\n".join(f"> {line}" for line in notes_text.splitlines())
        parts.append(f"**Speaker notes:**\n{quoted}")

    return "\n\n".join(parts)


class PPTXParser(BaseParser):
    """
    Parse PowerPoint presentations (.pptx, .ppt) using python-pptx.

    Each slide becomes a ## section with title, body text, tables,
    and speaker notes. Empty slides are skipped.
    """

    @property
    def supported_extensions(self) -> frozenset[str]:
        return frozenset({".pptx", ".ppt"})

    def parse(
        self,
        source_path: str,
        workspace: str,
        title: Optional[str] = None,
    ) -> DocumentNode:
        if not os.path.isfile(source_path):
            raise FileNotFoundError(f"PPTX file not found: {source_path!r}")

        try:
            from pptx import Presentation  # noqa: PLC0415
        except ImportError as e:
            raise ParserError(
                source_path,
                "python-pptx is not installed. Run: pip install 'python-pptx>=0.6.23'",
                original=e,
            )

        logger.info("Parsing PPTX: %s", source_path)

        try:
            prs = Presentation(source_path)
        except Exception as e:
            raise ParserError(
                source_path, f"python-pptx failed to open presentation: {e}", original=e
            )

        sections: list[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            md_section = _slide_to_markdown(slide, slide_num)
            if md_section:
                sections.append(md_section)
            else:
                logger.debug("Skipping empty slide %d in %s", slide_num, source_path)

        if not sections:
            raise ParserError(source_path, "Presentation contains no readable text content.")

        raw_text = "\n\n".join(sections)
        cleaned_text = clean(raw_text)

        return DocumentNode.from_file(
            source_path=source_path,
            text=cleaned_text,
            workspace=workspace,
            title=title,
        )
